import docx
import pptx

def extract_docx(file_path):
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text)
    return "\n".join(text)

def extract_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

with open("output.txt", "w", encoding="utf-8") as f:
    f.write("--- THESIS ---\n")
    f.write(extract_docx("e:\\PHYSICS Proj\\Quantum Teleportation Thesis.docx"))
    f.write("\n\n--- PPT ---\n")
    f.write(extract_pptx("e:\\PHYSICS Proj\\Physics_Project101.pptx"))
